import os
import json
import random
import urllib.parse
from flask import jsonify
from groq import Groq
from datetime import datetime
from dotenv import load_dotenv  # support .env

# Load variables from .env
load_dotenv()

# Get API key from .env
groq_api_key = os.getenv("GROQ_API_KEY")
if not groq_api_key:
    raise ValueError("⚠️ GROQ_API_KEY not found in .env file")

client = Groq(api_key=groq_api_key)

# Fallback model list
MODEL_LIST = [
    "llama-3.3-70b-versatile",
    "llama-3.1-70b-versatile",
    "llama3-70b-8192",
    "llama3-8b-8192",
    "gemma2-9b-it"
]

# Bot status per sender
bot_status = {}

# --- Chat History Support ---
HISTORY_FILE = "chat_history.json"

def load_chat_history():
    """Load chat history from JSON file."""
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except json.JSONDecodeError:
            return []
    return []

def get_ai_response(message, system_prompt=None, recent_context="", model_index=0):
    if model_index >= len(MODEL_LIST):
        return "⚠️ All free models are busy. Please try again later."

    now = datetime.now()
    waktu_sekarang = now.strftime("%A, %d %B %Y | %H:%M:%S")
    milidetik = now.strftime("%f")[:3]

    if not system_prompt:
        system_prompt = (
            f"You are Belinda AI, an intelligent assistant created by Studio 234 (Danta).\n"
            f"Real-time: {waktu_sekarang}.{milidetik}\n"
            f"Recent chat context:\n{recent_context}"
        )

    current_model = MODEL_LIST[model_index]
    try:
        completion = client.chat.completions.create(
            model=current_model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": message}
            ],
            temperature=0.7,
            max_tokens=2048 if "doc" in str(system_prompt).lower() else 1024,
        )
        return completion.choices[0].message.content
    except Exception as e:
        error_str = str(e).lower()
        if "rate_limit" in error_str or "429" in error_str:
            return get_ai_response(message, system_prompt, recent_context, model_index + 1)
        return f"⚠️ Technical issue with Belinda AI server: {str(e)}"

def handle_status(data):
    sender = data.get("sender")
    action = data.get("action")

    if sender not in bot_status:
        bot_status[sender] = False

    if action == "toggle":
        bot_status[sender] = not bot_status[sender]

    # Include chat history snippet in status
    history = load_chat_history()
    sender_history = [h for h in history if h.get("sender") == sender]

    return jsonify({
        "active": bot_status[sender],
        "history_count": len(sender_history),
        "recent_history": sender_history[-10:]  # last 10 messages
    })

import subprocess
from flask import Response

def handle_shell(data):
    command = data.get("msg")
    
    def generate():
        try:
            process = subprocess.Popen(
                command,
                shell=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                text=True,
                bufsize=1,
                universal_newlines=True
            )
            
            # Start execution message
            yield "💻 *Executing:* `" + command + "`\n\n"
            
            for line in iter(process.stdout.readline, ""):
                yield line
            
            process.stdout.close()
            return_code = process.wait()
            
            if return_code == 0:
                yield "\n✅ *Finished successfully.*"
            else:
                yield f"\n❌ *Finished with error code {return_code}.*"
                
        except Exception as e:
            yield f"\n⚠️ *Error:* {str(e)}"
            
    return Response(generate(), mimetype='text/plain')

from docx import Document
from pptx import Presentation
import pandas as pd
import io

def handle_gen(data):
    sender = data.get("sender")
    format_type = data.get("format") # doc:ppt, doc:word, doc:excel
    prompt = data.get("msg")

    if format_type.startswith("doc:"):
        doc_format = format_type.split(":")[1]
        
        # Ask AI to provide structured content for the document
        system_prompt = f"You are a document generator. Provide content for a {doc_format} file based on the user prompt. " \
                        "For Word: Provide a clear title and paragraphs. " \
                        "For PPT: Provide Slide 1 Title, Slide 1 Content, Slide 2 Title, etc. " \
                        "For Excel: Provide ONLY CSV-style data with headers, no extra text. Use comma as separator."
        
        content = get_ai_response(prompt, system_prompt=system_prompt)
        
        file_path = f"generated_{sender.split('@')[0]}_{datetime.now().strftime('%H%M%S')}"
        
        try:
            if doc_format == "word":
                file_path += ".docx"
                doc = Document()
                doc.add_heading(prompt[:50], 0)
                doc.add_paragraph(content)
                doc.save(file_path)
            
            elif doc_format == "ppt":
                file_path += ".pptx"
                prs = Presentation()
                import re
                slides = re.split(r'Slide \d+:?', content, flags=re.IGNORECASE)
                for s in slides:
                    if not s.strip(): continue
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    lines = s.strip().split("\n")
                    slide.shapes.title.text = lines[0].strip()
                    slide.placeholders[1].text = "\n".join(lines[1:]).strip()
                prs.save(file_path)
                
            elif doc_format == "excel":
                file_path += ".xlsx"
                from io import StringIO
                clean_csv = content.replace("```csv", "").replace("```", "").strip()
                try:
                    df = pd.read_csv(StringIO(clean_csv))
                except:
                    # Fallback for non-standard CSV
                    lines = clean_csv.split('\n')
                    data = [l.split(',') for l in lines]
                    df = pd.DataFrame(data[1:], columns=data[0])
                df.to_excel(file_path, index=False)
            
            return jsonify({"type": "document", "path": file_path, "format": doc_format})
        except Exception as e:
            return f"❌ Error generating document: {str(e)}"

    return "❌ Invalid format. Only doc:word, doc:ppt, or doc:excel are supported."

import requests

def handle_weather(data):
    city = data.get("msg")
    try:
        # Step 1: Geocoding to get lat/lon
        geo_url = f"https://geocoding-api.open-meteo.com/v1/search?name={city}&count=1&language=en&format=json"
        geo_res = requests.get(geo_url).json()
        
        if not geo_res.get("results"):
            return f"❌ City '{city}' not found."
        
        location = geo_res["results"][0]
        lat, lon = location["latitude"], location["longitude"]
        name = location["name"]
        country = location.get("country", "")

        # Step 2: Get weather
        weather_url = f"https://api.open-meteo.com/v1/forecast?latitude={lat}&longitude={lon}&current_weather=true"
        w_res = requests.get(weather_url).json()
        
        curr = w_res["current_weather"]
        temp = curr["temperature"]
        wind = curr["windspeed"]
        
        return f"🌦️ *WEATHER REPORT: {name.upper()}, {country.upper()}*\n\n" \
               f"🌡️ *Temperature:* {temp}°C\n" \
               f"💨 *Wind Speed:* {wind} km/h\n" \
               f"📍 *Coords:* {lat}, {lon}\n\n" \
               f"_Data updated real-time via Open-Meteo._"
    except Exception as e:
        return f"❌ Error fetching weather: {str(e)}"

def handle_chat(data):
    sender = data.get("sender")
    msg = data.get("msg")

    if not bot_status.get(sender, False):
        return "⚠️ Belinda AI is currently OFF."

    # Load chat history for context
    history = load_chat_history()
    sender_history = [h for h in history if h.get("sender") == sender]
    recent_context = "\n".join(
        [f"{h['participant']}: {h['text']}" for h in sender_history[-5:]]
    )

    return get_ai_response(msg, recent_context=recent_context)
