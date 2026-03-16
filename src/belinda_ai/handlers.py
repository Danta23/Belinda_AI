import os
import json
import random
import urllib.parse
from flask import jsonify
from groq import Groq
from datetime import datetime
from dotenv import load_dotenv  # support .env

# Load variables from .env
load_dotenv()

# Get API key from .env
groq_api_key = os.getenv("GROQ_API_KEY")
if not groq_api_key:
    raise ValueError("⚠️ GROQ_API_KEY not found in .env file")

client = Groq(api_key=groq_api_key)

# Fallback model list
MODEL_LIST = [
    "llama-3.3-70b-versatile",
    "llama-3.1-70b-versatile",
    "llama3-70b-8192",
    "llama3-8b-8192",
    "gemma2-9b-it"
]

# Bot status per sender
bot_status = {}

# --- Chat History Support ---
HISTORY_FILE = "chat_history.json"

def load_chat_history():
    """Load chat history from JSON file."""
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except (json.JSONDecodeError, IsADirectoryError):
            return []
    return []

def get_ai_response(message, system_prompt=None, recent_context="", model_index=0):
    if model_index >= len(MODEL_LIST):
        return "⚠️ All free models are busy. Please try again later."

    now = datetime.now()
    waktu_sekarang = now.strftime("%A, %d %B %Y | %H:%M:%S")
    milidetik = now.strftime("%f")[:3]

    ai_name = os.getenv("AI_NAME", "Belinda AI")
    ai_personality = os.getenv("AI_PERSONALITY", "an intelligent assistant created by Studio 234 (Danta)")
    ai_max_tokens = int(os.getenv("AI_MAX_TOKENS", 1024))

    if not system_prompt:
        system_prompt = (
            f"You are {ai_name}, {ai_personality}.\n"
            f"Real-time: {waktu_sekarang}.{milidetik}\n"
            f"Recent chat context:\n{recent_context}"
        )

    current_model = MODEL_LIST[model_index]
    try:
        completion = client.chat.completions.create(
            model=current_model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": message}
            ],
            temperature=0.7,
            max_tokens=2048 if "doc" in str(system_prompt).lower() else ai_max_tokens,
        )
        return completion.choices[0].message.content
    except Exception as e:
        error_str = str(e).lower()
        if "rate_limit" in error_str or "429" in error_str:
            return get_ai_response(message, system_prompt, recent_context, model_index + 1)
        return f"⚠️ Technical issue with Belinda AI server: {str(e)}"

def handle_status(data):
    sender = data.get("sender")
    action = data.get("action")

    if sender not in bot_status:
        bot_status[sender] = False

    if action == "toggle":
        bot_status[sender] = not bot_status[sender]

    # Include chat history snippet in status
    history = load_chat_history()
    sender_history = [h for h in history if h.get("sender") == sender]

    return jsonify({
        "active": bot_status[sender],
        "history_count": len(sender_history),
        "recent_history": sender_history[-10:]  # last 10 messages
    })

import subprocess
from flask import Response

def handle_shell(data):
    command = data.get("msg")
    
    def generate():
        try:
            process = subprocess.Popen(
                command,
                shell=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                text=True,
                bufsize=1,
                universal_newlines=True
            )
            
            # Start execution message
            yield "💻 *Executing:* `" + command + "`\n\n"
            
            for line in iter(process.stdout.readline, ""):
                yield line
            
            process.stdout.close()
            return_code = process.wait()
            
            if return_code == 0:
                yield "\n✅ *Finished successfully.*"
            else:
                yield f"\n❌ *Finished with error code {return_code}.*"
                
        except Exception as e:
            yield f"\n⚠️ *Error:* {str(e)}"
            
    return Response(generate(), mimetype='text/plain')

from docx import Document
from pptx import Presentation
import pandas as pd
import io

def handle_gen(data):
    sender = data.get("sender")
    format_type = data.get("format") # doc:ppt, doc:word, doc:excel, scr:ext, 3dm:ext
    prompt = data.get("msg")

    if format_type.startswith("doc:"):
        doc_format = format_type.split(":")[1]
        
        # Ask AI to provide structured content for the document
        system_prompt = f"You are a document generator. Provide content for a {doc_format} file based on the user prompt. " \
                        "For Word: Provide a clear title and paragraphs. " \
                        "For PPT: Provide Slide 1 Title, Slide 1 Content, Slide 2 Title, etc. " \
                        "For Excel: Provide ONLY CSV-style data with headers, no extra text. Use comma as separator."
        
        content = get_ai_response(prompt, system_prompt=system_prompt)
        
        file_path = f"generated_{sender.split('@')[0]}_{datetime.now().strftime('%H%M%S')}"
        
        try:
            if doc_format == "word":
                file_path += ".docx"
                doc = Document()
                doc.add_heading(prompt[:50], 0)
                doc.add_paragraph(content)
                doc.save(file_path)
            
            elif doc_format == "ppt":
                file_path += ".pptx"
                prs = Presentation()
                import re
                slides = re.split(r'Slide \d+:?', content, flags=re.IGNORECASE)
                for s in slides:
                    if not s.strip(): continue
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    lines = s.strip().split("\n")
                    slide.shapes.title.text = lines[0].strip()
                    slide.placeholders[1].text = "\n".join(lines[1:]).strip()
                prs.save(file_path)
                
            elif doc_format == "excel":
                file_path += ".xlsx"
                from io import StringIO
                clean_csv = content.replace("```csv", "").replace("```", "").strip()
                try:
                    df = pd.read_csv(StringIO(clean_csv))
                except:
                    # Fallback for non-standard CSV
                    lines = clean_csv.split('\n')
                    data = [l.split(',') for l in lines]
                    df = pd.DataFrame(data[1:], columns=data[0])
                df.to_excel(file_path, index=False)
            
            return jsonify({"type": "document", "path": file_path, "format": doc_format})
        except Exception as e:
            return f"❌ Error generating document: {str(e)}"

    elif format_type.startswith("scr:"):
        ext = format_type.split(":")[1].lower()
        
        # Supported programming language extensions
        valid_scr_exts = [
            "py", "lua", "pas", "js", "html", "css", "cpp", "c", "cs", "php", 
            "gd", "java", "go", "rs", "ts", "rb", "sql", "pl", "curl", "asm", 
            "vb", "ps1", "sh", "fish", "zsh", "bat", "vbs"
        ]
        
        if ext not in valid_scr_exts:
            return f"❌ Unsupported language extension: `.{ext}`. \n" \
                   f"Supported: {', '.join(valid_scr_exts)}."

        system_prompt = f"You are an expert programmer. Write a high-quality, complete, and commented {ext} script for: {prompt}. " \
                        "Provide ONLY the code, no conversational filler. Do not wrap in markdown code blocks unless necessary for clarity."
        
        code = get_ai_response(prompt, system_prompt=system_prompt)
        file_path = f"script_{sender.split('@')[0]}_{datetime.now().strftime('%H%M%S')}.{ext}"
        
        try:
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(code.replace("```" + ext, "").replace("```", "").strip())
            return jsonify({"type": "document", "path": file_path, "format": ext})
        except Exception as e:
            return f"❌ Error creating script: {str(e)}"

    elif format_type.startswith("3dm:"):
        ext = format_type.split(":")[1].lower()
        
        # Validation for 3D formats
        valid_3d_exts = ["fbx", "obj", "stl", "glb", "gltf"]
        if ext not in valid_3d_exts:
            return f"❌ Invalid 3D format: `.{ext}`. \nSupported: {', '.join(valid_3d_exts)}."

        from googlesearch import search as gsearch
        import shutil
        
        try:
            # --- STAGE 1: AGGRESSIVE SEARCH ---
            search_queries = [
                f'site:github.com "{prompt}" extension:{ext}',
                f'"{prompt}" filetype:{ext}',
                f'"{prompt}" 3d model direct link {ext}'
            ]
            
            for query in search_queries:
                try:
                    results = gsearch(query, num_results=10)
                    for url in results:
                        u_low = url.lower()
                        target_url = None
                        if "github.com" in u_low and f".{ext}" in u_low:
                            target_url = u_low.replace("github.com", "raw.githubusercontent.com").replace("/blob/", "/")
                        elif u_low.endswith(f".{ext}"):
                            target_url = url
                        
                        if target_url:
                            file_path = f"model_{sender.split('@')[0]}_{datetime.now().strftime('%H%M%S')}.{ext}"
                            resp = requests.get(target_url, stream=True, timeout=15)
                            if resp.status_code == 200 and "text/html" not in resp.headers.get("Content-Type", ""):
                                with open(file_path, 'wb') as f:
                                    for chunk in resp.iter_content(chunk_size=8192): f.write(chunk)
                                return jsonify({"type": "document", "path": file_path, "format": ext, "status": "found"})
                except: continue

            # --- STAGE 2: MULTI-SPACE AI GENERATION ---
            from gradio_client import Client
            # List of spaces to try in order of reliability
            spaces = ["huggingface-projects/Shap-E", "shinasum/Shap-E", "TencentARC/InstantMesh"]
            
            for space in spaces:
                try:
                    client = Client(space)
                    # Different spaces have different API names/params
                    if "InstantMesh" in space:
                        # InstantMesh usually needs an image, so we skip for direct text-to-3d
                        continue 
                    
                    result = client.predict(prompt, api_name="/text_to_3d") if "huggingface-projects" in space else \
                             client.predict(prompt, 1, 15, api_name="/predict")
                    
                    if result and os.path.exists(result):
                        ai_ext = "glb"
                        file_path = f"ai_gen_{sender.split('@')[0]}_{datetime.now().strftime('%H%M%S')}.{ai_ext}"
                        shutil.copy(result, file_path)
                        return jsonify({"type": "document", "path": file_path, "format": ai_ext, "status": "generated"})
                except Exception as e:
                    print(f"Space {space} failed: {e}")
                    continue

            # --- STAGE 3: SMART LINK FALLBACK ---
            return f"⚠️ *Status:* Search failed & AI servers are busy.\n\n" \
                   f"I couldn't generate the literal file right now, but you can download it here:\n" \
                   f"🔗 https://poly.pizza/search/{urllib.parse.quote(prompt)}\n" \
                   f"🔗 https://sketchfab.com/search?q={urllib.parse.quote(prompt)}"
                
        except Exception as e:
            return f"❌ 3D Engine Error: {str(e)}"

    return "❌ Invalid format. Supported: doc:word|ppt|excel, scr:ext, 3dm:ext."

import requests

def handle_weather(data):
    city = data.get("msg")
    try:
        # Step 1: Geocoding to get lat/lon
        geo_url = f"https://geocoding-api.open-meteo.com/v1/search?name={city}&count=1&language=en&format=json"
        geo_res = requests.get(geo_url).json()
        
        if not geo_res.get("results"):
            return f"❌ City '{city}' not found."
        
        location = geo_res["results"][0]
        lat, lon = location["latitude"], location["longitude"]
        name = location["name"]
        country = location.get("country", "")

        # Step 2: Get weather
        weather_url = f"https://api.open-meteo.com/v1/forecast?latitude={lat}&longitude={lon}&current_weather=true"
        w_res = requests.get(weather_url).json()
        
        curr = w_res["current_weather"]
        temp = curr["temperature"]
        wind = curr["windspeed"]
        
        return f"🌦️ *WEATHER REPORT: {name.upper()}, {country.upper()}*\n\n" \
               f"🌡️ *Temperature:* {temp}°C\n" \
               f"💨 *Wind Speed:* {wind} km/h\n" \
               f"📍 *Coords:* {lat}, {lon}\n\n" \
               f"_Data updated real-time via Open-Meteo._"
    except Exception as e:
        return f"❌ Error fetching weather: {str(e)}"

def handle_chat(data):
    sender = data.get("sender")
    msg = data.get("msg")

    if not bot_status.get(sender, False):
        return "⚠️ Belinda AI is currently OFF."

    # Load chat history for context
    history = load_chat_history()
    sender_history = [h for h in history if h.get("sender") == sender]
    recent_context = "\n".join(
        [f"{h['participant']}: {h['text']}" for h in sender_history[-5:]]
    )

    return get_ai_response(msg, recent_context=recent_context)

def handle_voice(req):
    sender = req.form.get("sender")
    
    if not bot_status.get(sender, False):
        return "⚠️ Belinda AI is currently OFF."

    if 'audio' not in req.files:
        return "❌ No audio file received."

    audio_file = req.files['audio']
    temp_path = f"temp_voice_{sender.split('@')[0]}_{datetime.now().strftime('%H%M%S')}.ogg"
    
    try:
        audio_file.save(temp_path)
        
        with open(temp_path, "rb") as file:
            transcription = client.audio.transcriptions.create(
                file=(temp_path, file.read()),
                model="whisper-large-v3",
                response_format="text",
                language="id" # Default id, whisper will still auto-detect if spoken in English generally.
            )
            
            transcript_text = transcription.strip()
            
            if not transcript_text:
                raise ValueError("Could not transcribe voice note.")
            
            # Now pass text to the AI model
            history = load_chat_history()
            sender_history = [h for h in history if h.get("sender") == sender]
            recent_context = "\n".join(
                [f"{h['participant']}: {h['text']}" for h in sender_history[-5:]]
            )
            
            # Optionally add a note to AI that this was spoken
            system_prompt = (
                f"You are Belinda AI, an intelligent assistant. The user just sent a voice note that was transcribed to text.\n"
                f"Recent chat context:\n{recent_context}"
            )
            
            ai_reply = get_ai_response(transcript_text, system_prompt=system_prompt, recent_context=recent_context)
            
            # Format the output nicely to show what the bot heard
            final_response = ai_reply
            
            return final_response
            
    except Exception as e:
        return f"❌ Error processing voice note: {str(e)}"
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)
